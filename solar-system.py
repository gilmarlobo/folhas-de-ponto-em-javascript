from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Paleta ──────────────────────────────────────────────
BG      = RGBColor(0x05, 0x05, 0x1A)   # azul-noite
GOLD    = RGBColor(0xFF, 0xD7, 0x00)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0xCC, 0xCC, 0xCC)

# ── Conteúdo ─────────────────────────────────────────────
bodies = {
    "Sol": {
        "tipo": "Estrela",
        "subtitulo": "O Coração do Sistema Solar",
        "fatos": [
            "Diâmetro: 1.392.700 km (109× a Terra)",
            "Massa: 1,989 × 10³⁰ kg (99,86% do sistema)",
            "Temperatura superficial: ~5.500 °C",
            "Temperatura do núcleo: ~15 milhões °C",
            "Composição: 73% hidrogênio, 25% hélio",
            "Distância ao centro da Via Láctea: ~26.000 anos-luz",
            "Rotação no equador: ~25 dias terrestres",
        ],
        "cor": "FFD700",
    },
    "Mercúrio": {
        "tipo": "Planeta Rochoso",
        "subtitulo": "O Mensageiro Veloz",
        "fatos": [
            "Diâmetro: 4.879 km — menor planeta do sistema",
            "Distância ao Sol: 57,9 milhões km (0,39 UA)",
            "Translação: apenas 88 dias terrestres",
            "Temperatura: −180 °C a +430 °C",
            "Sem atmosfera significativa",
            "Superfície repleta de crateras, similar à Lua",
            "Tem campo magnético fraco (~1% do da Terra)",
        ],
        "cor": "B5B5B5",
    },
    "Vênus": {
        "tipo": "Planeta Rochoso",
        "subtitulo": "A Estrela da Manhã",
        "fatos": [
            "Diâmetro: 12.104 km — quase igual à Terra",
            "Distância ao Sol: 108,2 milhões km (0,72 UA)",
            "Translação: 225 dias terrestres",
            "Temperatura média: 465 °C (o mais quente!)",
            "Atmosfera densa de CO₂ — efeito estufa extremo",
            "Rotação retrógrada: 1 dia venusiano = 243 dias terrestres",
            "Coberto por densas nuvens de ácido sulfúrico",
        ],
        "cor": "E8C97A",
    },
    "Terra": {
        "tipo": "Planeta Rochoso",
        "subtitulo": "O Pálido Ponto Azul",
        "fatos": [
            "Diâmetro: 12.742 km",
            "Distância ao Sol: 149,6 milhões km (1 UA)",
            "Translação: 365,25 dias",
            "Temperatura média: 15 °C",
            "Único planeta conhecido com vida",
            "71% da superfície coberta por água",
            "1 satélite natural: a Lua",
        ],
        "cor": "4FC3F7",
    },
    "Lua": {
        "tipo": "Satélite Natural",
        "subtitulo": "Nossa Companheira Eterna",
        "fatos": [
            "Diâmetro: 3.474 km",
            "Distância à Terra: ~384.400 km",
            "Translação ao redor da Terra: 27,3 dias",
            "Mesma face sempre voltada para a Terra (rotação síncrona)",
            "Temperatura: −173 °C a +127 °C",
            "Causa as marés terrestres",
            "12 humanos já pisaram na Lua (Apollo, 1969–1972)",
        ],
        "cor": "CCCCCC",
    },
    "Marte": {
        "tipo": "Planeta Rochoso",
        "subtitulo": "O Planeta Vermelho",
        "fatos": [
            "Diâmetro: 6.779 km",
            "Distância ao Sol: 227,9 milhões km (1,52 UA)",
            "Translação: 687 dias terrestres",
            "Temperatura média: −60 °C",
            "Possui o maior vulcão do sistema solar: Olympus Mons (21 km de altura)",
            "Dois pequenos satélites: Fobos e Deimos",
            "Missões ativas: Perseverance, Curiosity, InSight",
        ],
        "cor": "CF5C36",
    },
    "Júpiter": {
        "tipo": "Planeta Gasoso",
        "subtitulo": "O Rei dos Planetas",
        "fatos": [
            "Diâmetro: 139.820 km — maior do sistema solar",
            "Distância ao Sol: 778,5 milhões km (5,2 UA)",
            "Translação: 11,9 anos terrestres",
            "Grande Mancha Vermelha: tempestade há +300 anos",
            "95 satélites conhecidos, incluindo Io, Europa, Ganimedes e Calisto",
            "Rotação mais rápida: 1 dia = 9h 55min",
            "Campo magnético 20.000× mais forte que o da Terra",
        ],
        "cor": "C88B3A",
    },
    "Saturno": {
        "tipo": "Planeta Gasoso",
        "subtitulo": "O Senhor dos Anéis",
        "fatos": [
            "Diâmetro: 116.460 km",
            "Distância ao Sol: 1,43 bilhão km (9,5 UA)",
            "Translação: 29,5 anos terrestres",
            "Anéis formados de gelo e rocha, com até 282.000 km de diâmetro",
            "Densidade menor que a água — flutuaria num oceano!",
            "146 satélites — destaque para Titã com atmosfera densa",
            "Titã tem rios e lagos de metano líquido",
        ],
        "cor": "E4D191",
    },
    "Urano": {
        "tipo": "Gigante de Gelo",
        "subtitulo": "O Planeta Deitado",
        "fatos": [
            "Diâmetro: 50.724 km",
            "Distância ao Sol: 2,87 bilhões km (19,2 UA)",
            "Translação: 84 anos terrestres",
            "Eixo inclinado 97,77° — praticamente de lado!",
            "Temperatura mínima: −224 °C — mais frio que Netuno",
            "Possui 13 anéis tênues",
            "27 satélites com nomes de personagens de Shakespeare",
        ],
        "cor": "93E0E3",
    },
    "Netuno": {
        "tipo": "Gigante de Gelo",
        "subtitulo": "O Gigante Tempestuoso",
        "fatos": [
            "Diâmetro: 49.244 km",
            "Distância ao Sol: 4,5 bilhões km (30,1 UA)",
            "Translação: 165 anos terrestres",
            "Ventos mais rápidos do sistema: até 2.100 km/h",
            "Grande Mancha Escura — tempestade do tamanho da Terra",
            "14 satélites — destaque para Tritão (órbita retrógrada)",
            "Previsto matematicamente antes de ser observado (1846)",
        ],
        "cor": "4B70DD",
    },
    "Plutão": {
        "tipo": "Planeta Anão",
        "subtitulo": "O Exilado do Sistema Solar",
        "fatos": [
            "Diâmetro: 2.376 km — menor que a Lua",
            "Distância ao Sol: ~5,9 bilhões km (39,5 UA)",
            "Translação: 248 anos terrestres",
            "Rebaixado a planeta anão pela UAI em 2006",
            "Tem 5 satélites — Caronte é quase do mesmo tamanho",
            "Temperatura média: −225 °C",
            "New Horizons fez sobrevoo em julho de 2015",
        ],
        "cor": "A0785A",
    },
}

# ── Funções auxiliares ────────────────────────────────────

def set_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=24, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def add_rect_button(slide, text, left, top, width, height,
                    fill_color=RGBColor(0x1A, 0x1A, 0x4A),
                    text_color=WHITE, font_size=18, bold=True):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = GOLD
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    return shape


def add_hyperlink_to_shape(shape, slide_index, prs):
    """Adiciona hyperlink de navegação interna a um shape."""
    sp = shape._element
    # Cria relação de hyperlink
    rId = shape.part.relate_to(
        prs.slides[slide_index],
        'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide'
    )
    # Adiciona <a:hlinkClick> ao elemento <p:sp><p:nvSpPr><p:nvPr>
    nvSpPr = sp.find(qn('p:nvSpPr'))
    nvPr = nvSpPr.find(qn('p:nvPr'))
    if nvPr is None:
        nvPr = etree.SubElement(nvSpPr, qn('p:nvPr'))
    hlinkClick = etree.SubElement(nvPr, qn('a:hlinkClick'))
    hlinkClick.set(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id',
        rId
    )


# ── Criação dos slides ────────────────────────────────────

# Ordem dos slides
slide_names = ["capa", "menu", "sobre", "creditos"] + list(bodies.keys())
# índices: capa=0, menu=1, sobre=2, creditos=3, Sol=4, Mercúrio=5 … Plutão=14

slides = {}
slide_layout = prs.slide_layouts[6]  # em branco

for name in slide_names:
    sl = prs.slides.add_slide(slide_layout)
    set_bg(sl, BG)
    slides[name] = sl

idx = {name: i for i, name in enumerate(slide_names)}

# ════════════════════════════════════════════════
# SLIDE 0 — CAPA
# ════════════════════════════════════════════════
sl = slides["capa"]
add_textbox(sl, "🪐  SISTEMA SOLAR",
            Inches(1), Inches(1.5), Inches(11.33), Inches(1.5),
            font_size=54, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_textbox(sl, "Uma jornada pelo nosso vizinhança cósmica",
            Inches(1), Inches(3.2), Inches(11.33), Inches(0.8),
            font_size=22, color=LGRAY, align=PP_ALIGN.CENTER, italic=True)
add_textbox(sl, "Aula sobre Links, Zoom e Ações no PowerPoint",
            Inches(1), Inches(4.1), Inches(11.33), Inches(0.6),
            font_size=16, color=RGBColor(0x88, 0x88, 0xAA), align=PP_ALIGN.CENTER)

# Botão → Menu
btn = add_rect_button(sl, "▶  Iniciar", Inches(5.16), Inches(5.5), Inches(3), Inches(0.8))
add_hyperlink_to_shape(btn, idx["menu"], prs)

# ════════════════════════════════════════════════
# SLIDE 1 — MENU PRINCIPAL
# ════════════════════════════════════════════════
sl = slides["menu"]
add_textbox(sl, "MENU PRINCIPAL",
            Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8),
            font_size=36, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# Grade 4 colunas de planetas
planet_list = list(bodies.keys())
cols = 4
btn_w, btn_h = Inches(2.8), Inches(0.7)
start_x, start_y = Inches(0.4), Inches(1.4)
gap_x, gap_y = Inches(0.15), Inches(0.12)

for i, planet in enumerate(planet_list):
    col = i % cols
    row = i // cols
    x = start_x + col * (btn_w + gap_x)
    y = start_y + row * (btn_h + gap_y)
    color_hex = bodies[planet]["cor"]
    r, g, b = int(color_hex[0:2],16), int(color_hex[2:4],16), int(color_hex[4:6],16)
    fill = RGBColor(r, g, b)
    # escurece um pouco para fundo
    fill_dark = RGBColor(max(0,r-60), max(0,g-60), max(0,b-60))
    btn = add_rect_button(sl, planet, x, y, btn_w, btn_h,
                          fill_color=fill_dark, text_color=WHITE, font_size=16)
    add_hyperlink_to_shape(btn, idx[planet], prs)

# Botões de menu secundário
for label, target, xi in [("Sobre", "sobre", 2.5), ("Créditos", "creditos", 5.66), ("⬅ Capa", "capa", 8.83)]:
    btn = add_rect_button(sl, label,
                          Inches(xi), Inches(6.5), Inches(2.6), Inches(0.7),
                          fill_color=RGBColor(0x1A,0x1A,0x4A))
    add_hyperlink_to_shape(btn, idx[target], prs)

# ════════════════════════════════════════════════
# SLIDE 2 — SOBRE
# ════════════════════════════════════════════════
sl = slides["sobre"]
add_textbox(sl, "SOBRE ESTA APRESENTAÇÃO",
            Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8),
            font_size=34, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

sobre_txt = (
    "Esta apresentação foi criada como material didático para a aula de "
    "PowerPoint, com foco em três recursos avançados:\n\n"
    "🔗  LINKS INTERNOS — botões que levam a slides específicos sem seguir "
    "a ordem linear.\n\n"
    "🔍  ZOOM DE SLIDE — recurso nativo do PowerPoint que cria miniaturas "
    "clicáveis e animadas de outros slides.\n\n"
    "⚡  AÇÕES — gatilhos em shapes que executam hiperlinks, macros ou sons "
    "ao clicar ou passar o mouse.\n\n"
    "O tema escolhido — o Sistema Solar — inclui os 8 planetas oficiais, "
    "Plutão (planeta anão), o Sol e a Lua, totalizando 11 corpos celestes."
)
add_textbox(sl, sobre_txt,
            Inches(1), Inches(1.4), Inches(11.33), Inches(4.5),
            font_size=18, color=WHITE)

btn = add_rect_button(sl, "⬅ Voltar ao Menu",
                      Inches(5), Inches(6.5), Inches(3.33), Inches(0.7))
add_hyperlink_to_shape(btn, idx["menu"], prs)

# ════════════════════════════════════════════════
# SLIDE 3 — CRÉDITOS
# ════════════════════════════════════════════════
sl = slides["creditos"]
add_textbox(sl, "CRÉDITOS",
            Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8),
            font_size=40, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

cred_txt = (
    "📚  Conteúdo científico baseado em:\n"
    "  • NASA Solar System Exploration (solarsystem.nasa.gov)\n"
    "  • ESA — European Space Agency (esa.int)\n"
    "  • União Astronômica Internacional — IAU (iau.org)\n\n"
    "🎨  Design e diagramação: criado para fins educacionais.\n\n"
    "💡  Recursos de PowerPoint demonstrados:\n"
    "  • Hyperlinks internos (Inserir → Link → Este Documento)\n"
    "  • Zoom de Slide (Inserir → Zoom → Zoom de Slide)\n"
    "  • Ações (Inserir → Ação → Hiperlink para)\n\n"
    "© 2025 — Material Didático de Informática / PowerPoint Avançado"
)
add_textbox(sl, cred_txt,
            Inches(1.5), Inches(1.4), Inches(10.33), Inches(4.5),
            font_size=17, color=WHITE)

btn = add_rect_button(sl, "⬅ Voltar ao Menu",
                      Inches(5), Inches(6.5), Inches(3.33), Inches(0.7))
add_hyperlink_to_shape(btn, idx["menu"], prs)

# ════════════════════════════════════════════════
# SLIDES 4-14 — Corpos Celestes
# ════════════════════════════════════════════════
for planet, data in bodies.items():
    sl = slides[planet]
    color_hex = data["cor"]
    r,g,b = int(color_hex[0:2],16), int(color_hex[2:4],16), int(color_hex[4:6],16)
    accent = RGBColor(r, g, b)

    # Faixa colorida no topo
    top_bar = sl.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.08))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = accent
    top_bar.line.fill.background()

    # Círculo decorativo
    circle = sl.shapes.add_shape(
        9,  # oval
        Inches(10.5), Inches(0.5), Inches(2.5), Inches(2.5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent
    circle.line.fill.background()
    # Nome no círculo
    tf = circle.text_frame
    tf.text = planet[0]  # inicial
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(60)
    run.font.bold = True
    run.font.color.rgb = BG

    # Título
    add_textbox(sl, planet.upper(),
                Inches(0.5), Inches(0.2), Inches(9.5), Inches(0.9),
                font_size=42, bold=True, color=accent)
    add_textbox(sl, data["subtitulo"],
                Inches(0.5), Inches(1.05), Inches(9.5), Inches(0.5),
                font_size=18, color=LGRAY, italic=True)
    add_textbox(sl, data["tipo"].upper(),
                Inches(0.5), Inches(1.55), Inches(4), Inches(0.4),
                font_size=13, color=accent, bold=True)

    # Fatos
    fatos_txt = "\n".join(f"  •  {f}" for f in data["fatos"])
    add_textbox(sl, fatos_txt,
                Inches(0.5), Inches(2.1), Inches(9.5), Inches(4.5),
                font_size=16, color=WHITE)

    # Botão voltar
    btn = add_rect_button(sl, "⬅ Menu",
                          Inches(0.4), Inches(6.55), Inches(2), Inches(0.65),
                          fill_color=RGBColor(0x1A,0x1A,0x4A))
    add_hyperlink_to_shape(btn, idx["menu"], prs)

    # Navegação anterior / próximo
    planet_list_full = list(bodies.keys())
    pi = planet_list_full.index(planet)
    if pi > 0:
        prev = planet_list_full[pi - 1]
        btn2 = add_rect_button(sl, f"◀ {prev}",
                               Inches(2.7), Inches(6.55), Inches(2.5), Inches(0.65),
                               fill_color=RGBColor(0x1A,0x1A,0x4A))
        add_hyperlink_to_shape(btn2, idx[prev], prs)
    if pi < len(planet_list_full) - 1:
        nxt = planet_list_full[pi + 1]
        btn3 = add_rect_button(sl, f"{nxt} ▶",
                               Inches(9.83), Inches(6.55), Inches(2.5), Inches(0.65),
                               fill_color=RGBColor(0x1A,0x1A,0x4A))
        add_hyperlink_to_shape(btn3, idx[nxt], prs)

# ── Salvar ───────────────────────────────────────────────
output = "sistema_solar_links_zoom_acoes.pptx"
prs.save(output)
print(f"✅  Apresentação salva: {output}")
print(f"   Total de slides: {len(prs.slides)}")